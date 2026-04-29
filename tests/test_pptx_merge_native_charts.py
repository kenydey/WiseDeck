import zipfile
import asyncio
from io import BytesIO

from pptx import Presentation

from wisedeck.services.structured_export.schemas import (
    ChartConfigModel,
    ChartDataModel,
    ChartDatasetModel,
    StructuredSlideDeckModel,
    StructuredSlideModel,
)
from wisedeck.services.structured_export.service import merge_native_charts_into_pptx_bytes


def _blank_pptx_bytes(slide_count: int) -> bytes:
    prs = Presentation()
    for _ in range(slide_count):
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_merge_native_charts_adds_chart_xml():
    base_bytes = _blank_pptx_bytes(2)
    with zipfile.ZipFile(BytesIO(base_bytes), "r") as zf:
        names = zf.namelist()
    assert not [n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml")]

    deck = StructuredSlideDeckModel(
        title="Merge",
        language="zh",
        slides=[
            StructuredSlideModel(
                page_number=1,
                title="S1",
                chart_config=ChartConfigModel(
                    type="bar",
                    data=ChartDataModel(
                        labels=["A", "B"],
                        datasets=[ChartDatasetModel(label="Series", data=[1.0, 2.0])],
                    ),
                ),
            ),
            StructuredSlideModel(
                page_number=2,
                title="S2",
            ),
        ],
    )
    out_bytes = asyncio.run(merge_native_charts_into_pptx_bytes(base_bytes, deck=deck))
    assert out_bytes[:4] == b"PK\x03\x04"
    with zipfile.ZipFile(BytesIO(out_bytes), "r") as zf:
        names = zf.namelist()
    chart_parts = [n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml")]
    assert chart_parts, f"expected native chart xml in merged pptx, got sample: {names[:25]}"

