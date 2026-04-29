"""
Structured PPTX export: native python-pptx charts + Playwright homomorphic HTML (slides-html).
"""

from __future__ import annotations

import logging
import os
import tempfile
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.util import Inches

from wisedeck.services.structured_export.chart_mapper import deck_to_pptx_presentation_model
from wisedeck.services.structured_export._presenton.pptx_presentation_creator import (
    PptxPresentationCreator,
)
from wisedeck.services.structured_export.schemas import StructuredSlideDeckModel

logger = logging.getLogger(__name__)

def _pptx_base64_to_bytes(b64: str) -> bytes:
    import base64

    return base64.b64decode(b64.encode("ascii"))


async def export_dom_to_pptx_bytes_via_playwright(
    *,
    slides_html_url: str,
    timeout_ms: int = 180_000,
) -> bytes:
    """
    Run dom-to-pptx inside a headless Chromium (Playwright) against a same-origin slides-html document.

    Returns PPTX bytes. This is the server-side equivalent of the client export pipeline.
    """
    # IMPORTANT (Windows): Playwright requires asyncio subprocess support.
    # Some server event loops (uvicorn worker) can end up using a Selector loop that
    # raises NotImplementedError for create_subprocess_exec. To keep exports stable,
    # run Playwright in a dedicated thread with a fresh Proactor loop.
    from wisedeck.utils.thread_pool import run_blocking_io

    def _run() -> bytes:
        import asyncio
        import sys

        if sys.platform == "win32":
            proactor = getattr(asyncio, "WindowsProactorEventLoopPolicy", None)
            if proactor is not None:
                try:
                    asyncio.set_event_loop_policy(proactor())
                except Exception:
                    pass

        async def _async_run() -> bytes:
            from playwright.async_api import async_playwright

            # The dom-to-pptx runtime expects 1280x720 slide coordinates.
            viewport = {"width": 1280, "height": 720}

            runner_js = r"""
    async () => {
      const sleep = (ms) => new Promise((r) => setTimeout(r, ms));
      const waitWithTimeout = async (p, ms, msg) => {
        let t;
        try {
          return await Promise.race([
            p,
            new Promise((_, rej) => { t = setTimeout(() => rej(new Error(msg || 'timeout')), ms); })
          ]);
        } finally {
          if (t) clearTimeout(t);
        }
      };

      // Best-effort readiness gates: fonts + images.
      try {
        if (document.fonts && document.fonts.ready) {
          await waitWithTimeout(document.fonts.ready, 12000, 'fonts_ready_timeout');
        }
      } catch (_) {}
      try {
        const imgs = Array.from(document.images || []);
        await waitWithTimeout(Promise.all(imgs.map(img => {
          if (img.complete) return true;
          return new Promise((res) => {
            img.addEventListener('load', () => res(true), { once: true });
            img.addEventListener('error', () => res(false), { once: true });
          });
        })), 15000, 'images_ready_timeout');
      } catch (_) {}

      // Ensure dom-to-pptx runtime is available.
      const exporter = window.domToPptx;
      if (!exporter || typeof exporter.exportToPptx !== 'function') {
        throw new Error('domToPptx_not_loaded');
      }

      const pages = Array.from(document.querySelectorAll('.wisedeck-slide-page'));
      if (!pages.length) {
        throw new Error('no_slides_found');
      }

      // Force a stable 1280x720 viewport-like box for each slide root.
      for (const page of pages) {
        page.style.width = '1280px';
        page.style.height = '720px';
        page.style.minHeight = '720px';
        page.style.overflow = 'hidden';
      }

      async function* slideStream() {
        for (const page of pages) {
          const root = page.querySelector('.wisedeck-slide-root') || page;
          // Give the browser a frame to settle styles/layout.
          await sleep(0);
          yield root;
        }
      }

      const blob = await exporter.exportToPptx(slideStream(), {
        fileName: 'wisedeck.pptx',
        autoEmbedFonts: true,
        svgAsVector: false,
        skipDownload: true,
      });

      const buf = await blob.arrayBuffer();
      const u8 = new Uint8Array(buf);
      // Convert to base64 in chunks to avoid call stack limits.
      let binary = '';
      const chunk = 0x8000;
      for (let i = 0; i < u8.length; i += chunk) {
        binary += String.fromCharCode.apply(null, u8.subarray(i, i + chunk));
      }
      return btoa(binary);
    }
    """

            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                try:
                    context = await browser.new_context(viewport=viewport, ignore_https_errors=True)
                    page = await context.new_page()
                    await page.goto(slides_html_url, wait_until="domcontentloaded", timeout=timeout_ms)
                    # Extra readiness: wait for network to settle (best-effort).
                    try:
                        await page.wait_for_load_state("networkidle", timeout=15_000)
                    except Exception:
                        pass
                    # Load dom-to-pptx bundle from the same origin.
                    # (Keep the URL stable; cache-busting is optional for server-side.)
                    base_url = slides_html_url.split("/api/")[0].rstrip("/")
                    await page.add_script_tag(url=f"{base_url}/static/js/dom-to-pptx.bundle.js")
                    b64 = await page.evaluate(runner_js)
                    return _pptx_base64_to_bytes(str(b64))
                finally:
                    try:
                        await browser.close()
                    except Exception:
                        pass

        return asyncio.run(_async_run())

    return await run_blocking_io(_run)


async def screenshot_html_files_to_png_via_playwright(
    *,
    html_files: list[str],
    png_files: list[str],
    timeout_ms: int = 60_000,
) -> list[bool]:
    from wisedeck.utils.thread_pool import run_blocking_io

    def _run() -> list[bool]:
        import asyncio
        import sys
        from pathlib import Path

        if sys.platform == "win32":
            proactor = getattr(asyncio, "WindowsProactorEventLoopPolicy", None)
            if proactor is not None:
                try:
                    asyncio.set_event_loop_policy(proactor())
                except Exception:
                    pass

        async def _async_run() -> list[bool]:
            from playwright.async_api import async_playwright

            results: list[bool] = []
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                try:
                    context = await browser.new_context(
                        viewport={"width": 1280, "height": 720},
                        device_scale_factor=2,
                        ignore_https_errors=True,
                    )
                    page = await context.new_page()
                    for html_path, png_path in zip(html_files, png_files):
                        try:
                            url = Path(html_path).resolve().as_uri()
                            await page.goto(url, wait_until="domcontentloaded", timeout=timeout_ms)
                            try:
                                await page.wait_for_load_state("networkidle", timeout=8_000)
                            except Exception:
                                pass
                            await page.screenshot(
                                path=str(Path(png_path).resolve()),
                                type="png",
                                full_page=False,
                            )
                            results.append(True)
                        except Exception:
                            results.append(False)
                    try:
                        await page.close()
                    except Exception:
                        pass
                finally:
                    try:
                        await browser.close()
                    except Exception:
                        pass
            return results

        return asyncio.run(_async_run())

    return await run_blocking_io(_run)


async def build_pptx_bytes_from_slide_html_screenshots(
    *,
    slides: list[dict[str, Any]],
    export_base_url: str,
) -> bytes:
    """
    High-fidelity path (MVP): render each slide HTML to a PNG via Playwright, then build a PPTX with full-slide images.

    This is the minimal “same HTML as iframe” bridge: visuals match, but text is not editable.
    """
    tmp_dir = tempfile.mkdtemp(prefix="wd_html_pptx_")
    html_files: list[str] = []
    png_files: list[str] = []
    try:
        # Prepare per-slide HTML files.
        for i, slide in enumerate(slides):
            html_file = str(Path(tmp_dir) / f"slide_{i}.html")
            html_content = str(slide.get("html_content") or "")
            # Import locally to avoid widening import surface during tests.
            from wisedeck.web.route_modules.export_support import _prepare_html_for_file_based_export

            prepared = _prepare_html_for_file_based_export(html_content, export_base_url)
            Path(html_file).write_text(prepared, encoding="utf-8")
            html_files.append(html_file)

        # Screenshot each HTML file (server-side Playwright in a safe loop).
        for i, html_file in enumerate(html_files):
            png_file = str(Path(tmp_dir) / f"slide_{i}.png")
            png_files.append(png_file)

        ok_list = await screenshot_html_files_to_png_via_playwright(
            html_files=html_files,
            png_files=png_files,
        )
        for i, ok in enumerate(ok_list):
            if not ok:
                raise RuntimeError(f"Failed to screenshot slide {i + 1}")

        # Build PPTX with full-slide images (16:9).
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        for png_file in png_files:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                png_file,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        out = BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        # Best-effort cleanup.
        for p in png_files + html_files:
            try:
                Path(p).unlink(missing_ok=True)  # type: ignore[arg-type]
            except Exception:
                pass
        try:
            Path(tmp_dir).rmdir()
        except Exception:
            pass


async def merge_native_charts_into_pptx_bytes(
    base_pptx_bytes: bytes,
    *,
    deck: StructuredSlideDeckModel,
    chart_position_override_pt: tuple[int, int, int, int] | None = None,
) -> bytes:
    """
    Merge native editable charts into an existing .pptx produced by dom-to-pptx or other sources.

    MVP behavior: for each slide index that has chart_config, append a native chart shape to that slide.
    """
    base = BytesIO(base_pptx_bytes)
    prs = Presentation(base)

    # Keep deck -> slide index mapping simple and explicit: by order (0-based).
    for idx, slide_model in enumerate(deck.slides):
        if slide_model.chart_config is None:
            continue
        if idx >= len(prs.slides):
            # Ignore out-of-range deck slides; client export may be subset or mismatch.
            continue
        native_model = deck_to_pptx_presentation_model(
            StructuredSlideDeckModel(
                title=deck.title,
                language=deck.language,
                slides=[slide_model],
            )
        )
        # Expect: [title textbox, chart] shapes; best-effort find first chart model.
        chart_shape_model = None
        for s in (native_model.slides[0].shapes or []):
            if getattr(s, "shape_type", None) == "chart":
                chart_shape_model = s
                break
        if chart_shape_model is None:
            continue
        if chart_position_override_pt is not None:
            try:
                from wisedeck.services.structured_export._presenton.pptx_models import PptxPositionModel

                left, top, width, height = chart_position_override_pt
                chart_shape_model.position = PptxPositionModel(
                    left=int(left),
                    top=int(top),
                    width=int(width),
                    height=int(height),
                )
            except Exception:
                # Best-effort override; fall back to model-provided position.
                pass
        creator = PptxPresentationCreator(native_model, "")
        creator.add_chart(prs.slides[idx], chart_shape_model)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


async def build_pptx_bytes_from_deck(deck: StructuredSlideDeckModel) -> bytes:
    """Create .pptx bytes using vendored PptxPresentationCreator (editable charts)."""
    model = deck_to_pptx_presentation_model(deck)
    tmp = tempfile.mkdtemp(prefix="wd_struct_")
    try:
        creator = PptxPresentationCreator(model, tmp)
        await creator.create_ppt()
        out = Path(tmp) / f"{uuid.uuid4().hex}.pptx"
        creator.save(str(out))
        return out.read_bytes()
    finally:
        for f in Path(tmp).glob("*"):
            try:
                f.unlink()
            except OSError:
                pass
        try:
            Path(tmp).rmdir()
        except OSError:
            pass


def _structured_pptx_measurement_source() -> str:
    """
    Control which measurement/source pipeline `mode=auto` uses.

    Values:
    - homomorphic-editable (default): Playwright screenshots from slides-html + python-pptx chart merge
    - homomorphic: explicit screenshot pipeline + chart merge (same as editable naming for auto screenshots)
    - python-only: bypass screenshots; python-pptx only (lower fidelity)
    - wisedeck-html: reserved (currently maps to python-only or screenshots when inputs exist)

    ``render-service`` is ignored with a warning (maps to homomorphic-editable).
    """
    v = (os.environ.get("WISEDECK_STRUCTURED_PPTX_MEASUREMENT_SOURCE") or "").strip().lower()
    raw = v or "homomorphic-editable"
    if raw == "render-service":
        logger.warning(
            "WISEDECK_STRUCTURED_PPTX_MEASUREMENT_SOURCE=render-service is obsolete; using homomorphic-editable"
        )
        return "homomorphic-editable"
    return raw


async def export_structured_pptx_via_homomorphic_html(
    deck: StructuredSlideDeckModel,
    *,
    slides_for_same_html: list[dict[str, Any]] | None,
    export_base_url: str,
) -> bytes:
    """
    Homomorphic HTML path (WiseDeck iframe-equivalent):

    - Render `slides_for_same_html[i].html_content` via Playwright to PNGs
    - Build a PPTX with full-slide images (layout fidelity)
    - Merge python-pptx native editable charts (when chart_config exists)

    This intentionally prioritizes *visual parity* with the editor. Text/shape editability
    is not guaranteed (images), but charts become editable in PowerPoint.
    """
    if not slides_for_same_html:
        raise RuntimeError("Homomorphic export requires slides_data html_content (generate slides first).")
    base_pptx = await build_pptx_bytes_from_slide_html_screenshots(
        slides=slides_for_same_html,
        export_base_url=export_base_url,
    )
    # Default overlay region tuned for 16:9 full-slide backgrounds.
    # Units are points (pptx_models.PptxPositionModel uses Pt()).
    # Can be overridden via env var `WISEDECK_HOMOMORPHIC_CHART_BOX_PT=left,top,width,height`.
    override = None
    raw = (os.environ.get("WISEDECK_HOMOMORPHIC_CHART_BOX_PT") or "").strip()
    if raw:
        try:
            parts = [int(p.strip()) for p in raw.split(",") if p.strip()]
            if len(parts) == 4:
                override = (parts[0], parts[1], parts[2], parts[3])
        except Exception:
            override = None
    if override is None:
        override = (520, 170, 420, 300)
    return await merge_native_charts_into_pptx_bytes(
        base_pptx,
        deck=deck,
        chart_position_override_pt=override,
    )


async def export_structured_pptx_via_homomorphic_dom_to_pptx(
    deck: StructuredSlideDeckModel,
    *,
    project_id: str,
    export_base_url: str,
) -> bytes:
    """
    Phase-2 homomorphic export:
    - Server-side Playwright opens the same-origin hosted slides-html document
    - Runs dom-to-pptx to generate an editable base PPTX
    - Merges native editable charts (python-pptx)

    On failure, callers should fall back to `export_structured_pptx_via_homomorphic_html` (images).
    """
    slides_html_url = f"{export_base_url.rstrip('/')}/api/projects/{project_id}/internal/preview/slides-html"
    base_pptx = await export_dom_to_pptx_bytes_via_playwright(slides_html_url=slides_html_url)
    # Merge charts (overlay first; bbox alignment is Phase B).
    raw = (os.environ.get("WISEDECK_HOMOMORPHIC_CHART_BOX_PT") or "").strip()
    override = None
    if raw:
        try:
            parts = [int(p.strip()) for p in raw.split(",") if p.strip()]
            if len(parts) == 4:
                override = (parts[0], parts[1], parts[2], parts[3])
        except Exception:
            override = None
    if override is None:
        override = (520, 170, 420, 300)
    return await merge_native_charts_into_pptx_bytes(
        base_pptx,
        deck=deck,
        chart_position_override_pt=override,
    )


async def export_structured_pptx_auto(
    deck: StructuredSlideDeckModel,
    *,
    slides_for_same_html: Optional[list[dict[str, Any]]] = None,
    export_base_url: Optional[str] = None,
) -> bytes:
    """
    Default structured export for ``mode=auto``: driven by ``WISEDECK_STRUCTURED_PPTX_MEASUREMENT_SOURCE``
    (default ``homomorphic-editable``): Playwright screenshots from slides-html plus chart merge when applicable.

    Uses pure python-pptx when the measurement source is ``python-only`` or inputs are missing for screenshots.
    """
    source = _structured_pptx_measurement_source()
    if source in ("homomorphic-editable", "homomorphic_editable"):
        if not export_base_url:
            raise RuntimeError("Structured PPTX auto (homomorphic-editable) requires export_base_url")
        # Note: project_id is not available at this layer; callers should use explicit HTTP mode
        # for homomorphic_editable. Fall back to screenshot homomorphic when invoked from auto.
        if slides_for_same_html is not None:
            return await build_pptx_bytes_from_slide_html_screenshots(
                slides=slides_for_same_html,
                export_base_url=export_base_url,
            )
        return await build_pptx_bytes_from_deck(deck)
    if source in ("homomorphic", "homomorphic-html", "homomorphic_html"):
        if slides_for_same_html is None or not export_base_url:
            raise RuntimeError("Structured PPTX auto (homomorphic) requires slides_data and export_base_url")
        base_pptx = await build_pptx_bytes_from_slide_html_screenshots(
            slides=slides_for_same_html,
            export_base_url=export_base_url,
        )
        # Merge charts with default overlay region.
        raw = (os.environ.get("WISEDECK_HOMOMORPHIC_CHART_BOX_PT") or "").strip()
        override = None
        if raw:
            try:
                parts = [int(p.strip()) for p in raw.split(",") if p.strip()]
                if len(parts) == 4:
                    override = (parts[0], parts[1], parts[2], parts[3])
            except Exception:
                override = None
        if override is None:
            override = (520, 170, 420, 300)
        return await merge_native_charts_into_pptx_bytes(
            base_pptx,
            deck=deck,
            chart_position_override_pt=override,
        )
    if source in ("python-only", "python", "wisedeck-html"):
        if source == "wisedeck-html":
            if slides_for_same_html is not None and export_base_url:
                return await build_pptx_bytes_from_slide_html_screenshots(
                    slides=slides_for_same_html,
                    export_base_url=export_base_url,
                )
            logger.warning(
                "WISEDECK_STRUCTURED_PPTX_MEASUREMENT_SOURCE=wisedeck-html missing slides/base_url; using python-only export"
            )
        return await build_pptx_bytes_from_deck(deck)

    logger.warning("Unknown WISEDECK_STRUCTURED_PPTX_MEASUREMENT_SOURCE=%s; using python-only", source)
    return await build_pptx_bytes_from_deck(deck)
